"""Models to manage grahical attributes of presentation layouts, paragraphs, and charts."""

# Third-Party Libraries
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION, XL_LEGEND_POSITION, XL_TICK_MARK
from pptx.util import Pt


class Paragraph:
    """Simple class to call text frame attributes."""

    def shapes(self, slide):
        """Define a text frame."""
        shape = None

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
        return shape

    def shapes_find(self, slide):
        """Find text frames."""
        for shape in slide.shapes:
            print(shape.name)
        return

    def text_frame(self, shape):
        """Clear content of a text frame."""
        text_frame = shape.text_frame
        text_frame.clear()
        return text_frame

    def text_frame_ov_val(self, slide, shape, name):
        """Get overview page values."""
        text_frame = shape.text_frame
        for shape in slide.shapes:
            if shape.name == name:
                text_frame = shape.text_frame
                text_frame.clear()
        return text_frame

    def text_frame_key_metric(self, slide, shape, name):
        """Get metrics text."""
        text_frame = shape.text_frame
        for shape in slide.shapes:
            if shape.name == name:
                text_frame = shape.text_frame
                text_frame.clear()
        return text_frame

    def text_style_title(self, font):
        """Text style for cover page title."""
        font.name = "Calibri"
        font.size = Pt(28)
        font.color.rgb = RGBColor(255, 255, 255)
        return font

    def text_style_key_metric(self, font):
        """Text style for metrics."""
        font.name = "Calibri"
        font.size = Pt(12)
        font.color.rgb = RGBColor(255, 255, 255)
        return font

    def text_style_summary(self, font):
        """Text style for metrics."""
        font.name = "Calibri"
        font.size = Pt(10)
        font.color.rgb = RGBColor(0, 0, 0)
        return font

    def text_style_ov_val(self, font):
        """Text style for overview page values."""
        font.name = "Calibri"
        font.size = Pt(28)
        font.color.rgb = RGBColor(3, 37, 126)
        return font


class Graph:
    """Simple class to call chart attributes."""

    def bar_domain(self, xAxisTitle, yAxisTitle, chart):
        """Medium bar chart."""
        chart.font.size = Pt(10)
        chart.font.rgb = (20, 200, 100)
        chart.has_title = False
        x_axis = chart.category_axis.axis_title.text_frame
        x_axis.text = xAxisTitle
        x_axis.paragraphs[0].font.bold = False
        y_axis = chart.value_axis.axis_title.text_frame
        y_axis.text = yAxisTitle
        y_axis.paragraphs[0].font.bold = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(8)
        data_labels.font.color.rgb = RGBColor(211, 211, 211)
        data_labels.position = XL_LABEL_POSITION.INSIDE_END
        return

    def bar_sm(self, slide, chart):
        """Small bar chart."""
        chart.has_title = False
        chart.font.size = Pt(10)
        chart.font.rgb = (20, 200, 100)
        chart.value_axis.visible = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(8)
        data_labels.font.color.rgb = RGBColor(211, 211, 211)
        data_labels.position = XL_LABEL_POSITION.INSIDE_END
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.format.line.width = 0
        return

    def bar(self, slide, chart):
        """Medium bar chart."""
        chart.font.size = Pt(10)
        chart.font.rgb = (20, 200, 100)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels

        data_labels.font.size = Pt(7)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_LABEL_POSITION.ABOVE
        return

    def bar_med_100(self, slide, chart):
        """Medium bar chart with percentages."""
        chart.font.size = Pt(10)
        chart.font.rgb = (20, 200, 100)

        plot = chart.plots[0]
        plot.has_data_labels = True

        category_axis = chart.category_axis
        category_axis.has_major_gridlines = False
        category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
        category_axis.tick_labels.font.italic = True
        category_axis.tick_labels.font.size = Pt(8)

        data_labels = plot.data_labels
        data_labels.font.size = Pt(8)
        data_labels.font.color.rgb = RGBColor(211, 211, 211)
        data_labels.position = XL_LABEL_POSITION.INSIDE_END

        value_axis = chart.value_axis
        value_axis.maximum_scale = 100.0
        value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
        value_axis.has_minor_gridlines = True
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False

        tick_labels = value_axis.tick_labels
        tick_labels.number_format = '0"%"'
        tick_labels.font.bold = False
        tick_labels.font.size = Pt(8)

        return

    def line_med(self, slide, chart):
        """Medium line chart."""
        chart.font.size = Pt(10)
        chart.font.rgb = (20, 200, 100)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels

        data_labels.font.size = Pt(7)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
        data_labels.position = XL_LABEL_POSITION.ABOVE
        return


class Table:
    """Simple class to create tables."""

    def summary_table(self, table, rows, cols, df, merge):
        """Create summery table attributes."""
        if merge:
            for row in range(rows):
                table.cell(row, 0).merge(table.cell(row, 1))
        for col in cols:
            for row in range(rows):
                cell = table.cell(row, col)
                cell.fill.solid()
                if row != 0:
                    if merge and col == 2:
                        cell.text = str(df.iat[row - 1, col - 1])
                    else:
                        cell.text = str(df.iat[row - 1, col])
                    cell.fill.fore_color.rgb = RGBColor(231, 230, 230)
                else:
                    cell.fill.fore_color.rgb = RGBColor(180, 197, 231)
                cell.text_frame.paragraphs[0].font.size = Pt(9)
                cell.text_frame.paragraphs[0].font.name = "Calibri"
                cell.text_frame.paragraphs[0].font.bold = False
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
