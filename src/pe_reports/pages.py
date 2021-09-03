"""Collect and distribute graphical data to readable charts in the presentation."""

# Standard Python Libraries
import logging

# Third-Party Libraries
import pandas as pd
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from .stylesheet import Graph, Paragraph

# Use the following function to locate "Text Box" element in slide.
# find_shape = Paragraph.shapes_find(prs, slide)


class Pages:
    """Class containing the presentation page information."""

    @staticmethod
    def add_overview_value(run, df):
        """Add summary stats to the overview page."""
        run.text = str(int(df.iloc[0]["count"]))
        font = run.font
        Paragraph.text_style_ov_val(font)

    @staticmethod
    def insert_chart(slide, df_loc, col_names, chart_type, size, x, y, cx, cy, stacked):
        """Insert charts into the PowerPoint."""
        chart = CategoryChartData()
        try:
            df = pd.read_csv(df_loc).fillna(0)
            chart.categories = list(df.columns.values)
            chart.add_series(col_names[0], list(df.loc[0]))
            if stacked:
                chart.add_series(col_names[1], list(df.loc[1]))

            # TODO: Remove hard-coded graph size and positioning values
            # Issue 9: https://github.com/cisagov/pe-reports/issues/9
            chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart).chart
            if size == "small":
                Graph.chart_sm(chart)
            elif size == "medium":
                Graph.chart_med(chart)

        except FileNotFoundError as not_found:
            logging.error("%s : There is no customer data.", not_found)

    @staticmethod
    def cover(prs):
        """Page 1: Cover page of presentation."""
        slide = prs.slides[0]
        shape = Paragraph.shapes(slide)
        text_frame = Paragraph.text_frame(shape)
        frame = text_frame.paragraphs[0]
        run = frame.add_run()
        try:
            # TODO: Remove hard-coded file locations
            # Issue 8: https://github.com/cisagov/pe-reports/issues/8
            df_customer = pd.read_csv("src/pe_reports/data/csv/dhs_customer.csv")
            dates = (
                df_customer.iloc[0]["start_date"]
                + " to "
                + df_customer.iloc[0]["end_date"]
            )
            run.text = (
                "Prepared for: "
                + str(df_customer.iloc[0]["name"])
                + "\nReporting period: "
                + str(dates)
            )
            font = run.font
            Paragraph.text_style_title(font)
        except FileNotFoundError as not_found:
            logging.error("%s : There is no customer data.", not_found)

        return prs

    @staticmethod
    def overview(prs):
        """Page 2: Posture & Exposure Report Overview."""
        slide = prs.slides[1]

        # Overview Value - Credentials exposed in recent posts
        ov_val_01 = Paragraph.text_frame_ov_val(slide, "TextBox 2")
        frame = ov_val_01.paragraphs[0]
        run = frame.add_run()
        try:
            # TODO: Remove hard-coded file locations
            # Issue 8: https://github.com/cisagov/pe-reports/issues/8
            df_creds = pd.read_csv("src/pe_reports/data/csv/dhs_creds.csv")

            Pages.add_overview_value(run, df_creds)
        except FileNotFoundError as not_found:
            logging.error("%s : There is no customer data.", not_found)

        # Overview Value - Suspected domain masquerading alerts
        ov_val_02 = Paragraph.text_frame_ov_val(slide, "TextBox 82")
        frame = ov_val_02.paragraphs[0]
        run = frame.add_run()
        try:
            # TODO: Remove hard-coded file locations
            # Issue 8: https://github.com/cisagov/pe-reports/issues/8
            df_domains = pd.read_csv("src/pe_reports/data/csv/dhs_domains.csv")

            Pages.add_overview_value(run, df_domains)
        except FileNotFoundError as not_found:
            logging.error("%s : There is no customer data.", not_found)

        # Overview Value - Active malware associations
        ov_val_03 = Paragraph.text_frame_ov_val(slide, "TextBox 86")
        frame = ov_val_03.paragraphs[0]
        run = frame.add_run()
        try:
            # TODO: Remove hard-coded file locations
            # Issue 8: https://github.com/cisagov/pe-reports/issues/8
            df_malware = pd.read_csv("src/pe_reports/data/csv/dhs_malware.csv")

            Pages.add_overview_value(run, df_malware)
        except FileNotFoundError as not_found:
            logging.error("%s : There is no customer data.", not_found)

        # Overview Value - Web and dark web mentions
        ov_val_04 = Paragraph.text_frame_ov_val(slide, "TextBox 87")
        frame = ov_val_04.paragraphs[0]
        run = frame.add_run()
        try:
            # TODO: Remove hard-coded file locations
            # Issue 8: https://github.com/cisagov/pe-reports/issues/8
            df_vulns = pd.read_csv("src/pe_reports/data/csv/dhs_vulns.csv")
            Pages.add_overview_value(run, df_vulns)
        except FileNotFoundError as not_found:
            logging.error("%s : There is no customer data.", not_found)

        # Overview Value - Credentials exposed in recent posts
        ov_val_05 = Paragraph.text_frame_ov_val(slide, "TextBox 90")
        frame = ov_val_05.paragraphs[0]
        run = frame.add_run()
        try:
            # TODO: Remove hard-coded file locations
            # Issue 45: https://github.com/cisagov/pe-reports/issues/45
            df_web = pd.read_csv("src/pe_reports/data/csv/dhs_web.csv")
            Pages.add_overview_value(run, df_web)
        except FileNotFoundError as not_found:
            logging.error("%s : There is no customer data.", not_found)

        # Bar Graph - Top Level Domains used for masquerading
        # TODO: Remove hard-coded file locations
        # Issue 45: https://github.com/cisagov/pe-reports/issues/45
        df_loc = "src/pe_reports/data/csv/dhs_tld_df.csv"
        col_names = ["Top Level Domains"]
        chart_type = XL_CHART_TYPE.COLUMN_STACKED
        size = "small"
        x, y, cx, cy = Inches(3.5), Inches(2.9), Inches(2), Inches(1.5)
        stacked = False
        Pages.insert_chart(
            slide, df_loc, col_names, chart_type, size, x, y, cx, cy, stacked
        )

        # Bar Graph - Sources of credential exposures
        # TODO: Remove hard-coded file locations
        # Issue 45: https://github.com/cisagov/pe-reports/issues/45
        df_loc = "src/pe_reports/data/csv/dhs_ce_df.csv"
        col_names = ["Top Level Domains"]
        chart_type = XL_CHART_TYPE.COLUMN_STACKED
        size = "small"
        x, y, cx, cy = Inches(5.75), Inches(2.9), Inches(2), Inches(1.5)
        stacked = False
        Pages.insert_chart(
            slide, df_loc, col_names, chart_type, size, x, y, cx, cy, stacked
        )

        # Line Graph - Sources of credential exposures
        # TODO: Remove hard-coded file locations
        # Issue 45: https://github.com/cisagov/pe-reports/issues/45
        df_loc = "src/pe_reports/data/csv/dhs_web_df.csv"
        col_names = ["Web", "Dark Web"]
        chart_type = XL_CHART_TYPE.LINE
        size = "medium"
        x, y, cx, cy = Inches(3.5), Inches(5), Inches(4.8), Inches(2.15)
        stacked = True
        Pages.insert_chart(
            slide, df_loc, col_names, chart_type, size, x, y, cx, cy, stacked
        )

        # Bar Graph - Active malware associations
        # TODO: Remove hard-coded file locations
        # Issue 45: https://github.com/cisagov/pe-reports/issues/45
        df_loc = "src/pe_reports/data/csv/dhs_ma_df.csv"
        col_names = ["Sources"]
        chart_type = XL_CHART_TYPE.COLUMN_STACKED_100
        size = "small"
        x, y, cx, cy = Inches(8.25), Inches(2.9), Inches(4.5), Inches(2.0)
        stacked = False
        Pages.insert_chart(
            slide, df_loc, col_names, chart_type, size, x, y, cx, cy, stacked
        )

        # Bar Graph - inferred vulnerabilities found via external observation
        # TODO: Remove hard-coded file locations
        # Issue 45: https://github.com/cisagov/pe-reports/issues/45
        df_loc = "src/pe_reports/data/csv/dhs_iv_df.csv"
        col_names = ["Sources"]
        chart_type = XL_CHART_TYPE.COLUMN_STACKED_100
        size = "small"
        x, y, cx, cy = Inches(8.25), Inches(4.9), Inches(4.5), Inches(2.0)
        stacked = False
        Pages.insert_chart(
            slide, df_loc, col_names, chart_type, size, x, y, cx, cy, stacked
        )

        return prs
