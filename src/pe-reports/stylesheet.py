"""Need to insert docstring here."""
# Third-Party Libraries
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.util import Pt


class Paragraph:
    """Need to insert docstring here."""

    def shapes(self, slide):
        """Need to insert docstring here."""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
        return shape

    def shapes_find(self, slide):
        """Need to insert docstring here."""
        for shape in slide.shapes:
            print(shape.name)
        return

    def text_frame(self, shape):
        """Need to insert docstring here."""
        text_frame = shape.text_frame
        text_frame.clear()
        return text_frame

    def text_frame_ov_val(self, slide, shape, name):
        """Need to insert docstring here."""
        for shape in slide.shapes:
            if shape.name == name:
                text_frame = shape.text_frame
                text_frame.clear()
        return text_frame

    def text_frame_key_metric(self, slide, shape, name):
        """Need to insert docstring here."""
        for shape in slide.shapes:
            if shape.name == name:
                text_frame = shape.text_frame
                text_frame.clear()
        return text_frame

    def text_style_title(self, font):
        """Need to insert docstring here."""
        font.name = "Calibri"
        font.size = Pt(28)
        font.color.rgb = RGBColor(255, 255, 255)
        return font

    def text_style_key_metric(self, font):
        """Need to insert docstring here."""
        font.name = "Calibri"
        font.size = Pt(12)
        font.color.rgb = RGBColor(255, 255, 255)
        return font

    def text_style_ov_val(self, font):
        """Need to insert docstring here."""
        font.name = "Calibri"
        font.size = Pt(28)
        font.color.rgb = RGBColor(3, 37, 126)
        return font


class Graph:
    """Need to insert docstring here."""

    def bar(self, slide, chart):
        """Need to insert docstring here."""
        chart.font.size = Pt(10)
        chart.font.rgb = (20, 200, 100)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        return

    def bar_sm(self, slide, chart):
        """Need to insert docstring here."""
        chart.font.size = Pt(10)
        chart.font.rgb = (20, 200, 100)
        return

    def bar_med_100(self, slide, chart):
        """Need to insert docstring here."""
        chart.font.size = Pt(10)
        chart.font.rgb = (20, 200, 100)
        return

    def line_med(self, slide, chart):
        """Need to insert docstring here."""
        chart.font.size = Pt(10)
        chart.font.rgb = (20, 200, 100)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        return
